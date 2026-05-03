"""Generate presentation slides for Tree-HyperLISTA paper."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE = RGBColor(0x1A, 0x47, 0x8A)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFA)
GREEN = RGBColor(0x1D, 0x8D, 0x48)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
FIG = "tree_figures"


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, l, t, w, h, txt, sz=18, b=False, c=DARK, a=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.bold = b
    p.font.color.rgb = c
    p.font.name = "Calibri"
    p.alignment = a
    return tf


def bl(slide, l, t, w, h, items, sz=18, c=DARK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def bar(slide, l, t, w=0.06, h=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def img(slide, path, l, t, w=None, h=None):
    fp = os.path.join(FIG, path) if not os.path.isabs(path) else path
    if not os.path.exists(fp):
        tb(slide, l, t, 4, 1, f"[Missing: {path}]", 14, c=RED)
        return
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    slide.shapes.add_picture(fp, Inches(l), Inches(t), **kw)


def sec(slide, title):
    bg(slide, BLUE)
    tb(slide, 1, 2.5, 11, 2, title, 44, b=True, c=WHITE, a=PP_ALIGN.CENTER)


def title_bar(slide, title):
    bg(slide)
    tb(slide, 0.8, 0.4, 11, 0.8, title, 36, b=True, c=BLUE)
    bar(slide, 0.8, 1.1, h=0.06, w=3.0)


def make_table(slide, header, rows, l, t, w, h):
    ts = slide.shapes.add_table(len(rows)+1, len(header), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = ts.table
    for j, hdr in enumerate(header):
        cell = tbl.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    for i, row in enumerate(rows):
        ours = "TH-" in row[0] or "Tree-HyperLISTA" in row[0]
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = ours
                p.font.color.rgb = GREEN if ours else DARK
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8,0xF8,0xE8) if ours else (RGBColor(0xF8,0xF8,0xF8) if i%2==0 else WHITE)


# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_BG)
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()
tb(s, 1, 1.8, 11.3, 1.5, "Tree-HyperLISTA", 52, b=True, c=BLUE, a=PP_ALIGN.CENTER)
tb(s, 1, 3.0, 11.3, 1, "Ultra-Lightweight Deep Unfolding for Tree-Sparse Recovery", 28, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 1, 4.3, 11.3, 0.7, "Shresth Verma    Aditya Agrawal", 22, b=True, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 1, 4.9, 11.3, 0.7, "Department of CSE, IIT Bombay", 18, c=GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 2: Motivation ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Why Tree-Sparse Recovery?")
bl(s, 1, 1.6, 5.8, 5, [
    "Sparse recovery: reconstruct x from y = Ax + e",
    "Tree-structured sparsity arises naturally in:",
    "   \u2022 Wavelet decompositions of images",
    "   \u2022 Hierarchical signal representations",
    "   \u2022 Multi-resolution analysis",
    "",
    "Key constraint: if a child is active,",
    "its parent must also be active",
], 20)
bl(s, 7, 1.6, 5.5, 5, [
    "Gap in existing methods:",
    "",
    "Classical (ISTA, FISTA): no learning \u2192 slow",
    "LISTA (1.56M params): ignores tree structure",
    "HyperLISTA (3 params): elementwise only",
    "",
    "\u2192 No method combines tree structure",
    "   with lightweight deep unfolding",
], 20)

# === SLIDE 3: Method ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Tree-HyperLISTA: Method")
bl(s, 1, 1.5, 6, 5.5, [
    "Only 3 learned hyperparameters (c1, c2, c3)",
    "",
    "Each unfolded layer performs:",
    "  1. Momentum step:  z = x + \u03b2(x \u2212 x_prev)",
    "  2. Gradient step:  u = z + W\u1d40(y \u2212 Az)",
    "  3. Tree scoring:   si = |ui| + \u03c1 \u03a3 s_child",
    "  4. Support select: S = TopKTree(s, K)",
    "  5. Tree threshold: soft-thresh within S",
    "",
    "Per-layer \u03b2, \u03b8, K derived from (c1,c2,c3)",
    "via adaptive formulas \u2014 NOT learned per-layer",
], 19)
bl(s, 7.5, 1.5, 5, 5, [
    "Three support selection modes:",
    "",
    "M1 \u2014 Hard tree projection:",
    "    greedy top-K with ancestor closure",
    "",
    "M2 \u2014 Threshold + closure:",
    "    threshold then enforce tree consistency",
    "",
    "M3 \u2014 Hybrid (recommended):",
    "    top-K tree projection + soft threshold",
    "    \u2192 best NMSE with lowest variance",
], 19)

# === SLIDE 4: Key Innovation ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Key Innovation: Tree-Aware Proximal Operator")
bl(s, 1, 1.6, 5.8, 5.5, [
    "Standard (LISTA/HyperLISTA):",
    "   xi = sign(ui) max(|ui| \u2212 \u03b8, 0)",
    "   \u2192 treats each coefficient independently",
    "   \u2192 ignores tree structure completely",
    "",
    "Tree-HyperLISTA:",
    "   1. Score each node using subtree info",
    "   2. Select tree-consistent support S",
    "   3. Threshold only within S",
    "   \u2192 exploits parent-child dependencies",
    "   \u2192 O(n) via bottom-up aggregation",
], 20)
bl(s, 7.2, 1.6, 5.5, 5.5, [
    "What makes it \"ultra-lightweight\":",
    "",
    "\u2022 Same 3 global hyperparameters as HyperLISTA",
    "\u2022 W matrix is precomputed (not learned)",
    "\u2022 Tree topology is fixed (not learned)",
    "\u2022 \u03c1 decay is fixed at 0.5 (not learned)",
    "",
    "The ONLY change vs HyperLISTA is",
    "replacing soft-threshold with tree operator",
    "",
    "\u2192 Structure in the operator, not parameters",
], 20)

# === SLIDE 5: Section header ===
s = prs.slides.add_slide(prs.slide_layouts[6])
sec(s, "Experimental Results")

# === SLIDE 6: Core Results ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Core Results: Synthetic Tree-Sparse Recovery")
tb(s, 1, 1.4, 11, 0.5, "n=255, m=127, K=16 layers, 30 dB SNR, 5 seeds", 16, c=GRAY)
make_table(s,
    ["Model", "NMSE (dB)", "Precision", "Recall", "#Params"],
    [
        ["Tree-ISTA",       "\u221206.00 \u00b1 0.02",  "0.759", "0.734", "0"],
        ["Tree-FISTA",      "\u221208.34 \u00b1 0.02",  "0.840", "0.816", "0"],
        ["LISTA",           "\u221213.18 \u00b1 0.08", "0.255", "0.997", "1.56M"],
        ["HyperLISTA",      "\u221215.91 \u00b1 6.91", "0.791", "0.967", "3"],
        ["TH-Hard (M1)",    "\u221228.80 \u00b1 0.11", "0.740", "1.000", "3"],
        ["TH-Hybrid (M3)",  "\u221228.66 \u00b1 0.10", "0.663", "1.000", "3"],
    ], 1, 2, 11, 3.5)
bl(s, 1, 5.8, 11, 1.2, [
    "TH-Hybrid: +12.8 dB over HyperLISTA (same 3 params)  |  +15.5 dB over LISTA (1.56M params)  |  Perfect recall",
], 18, c=BLUE)

# === SLIDE 7: Convergence ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Convergence & Support Recovery")
img(s, "tree_fig1_convergence.png", 0.5, 1.5, w=6)
img(s, "tree_fig2_support.png", 6.8, 1.5, w=6)
tb(s, 0.5, 6.2, 12, 0.8, "Tree-HyperLISTA converges to \u221228.66 dB in 16 layers with perfect support recall", 18, c=GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 8: Mismatch ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Robustness: Mismatch Experiments")
img(s, "tree_fig3_snr_db.png", 0.3, 1.4, w=6.2)
img(s, "tree_fig4_target_sparsity.png", 6.8, 1.4, w=6.2)
tb(s, 0.5, 6, 12, 1, "Tree-HyperLISTA degrades gracefully under SNR and sparsity mismatch", 18, c=GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 9: Ablation ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Ablation: Sensitivity & Mechanisms")
img(s, "tree_fig6_sensitivity.png", 0.3, 1.4, w=6.2)
img(s, "tree_fig7_mechanisms.png", 6.8, 1.4, w=6.2)
bl(s, 0.5, 5.8, 12, 1.5, [
    "Sensitivity: broad optimal basin for c2 (momentum); c1 (threshold) most critical",
    "Mechanisms: M1 (hard) and M3 (hybrid) comparable; M2 (threshold) unstable",
], 17, c=GRAY)

# === SLIDE 10: Image CS ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Wavelet-Domain Image Compressed Sensing")
img(s, "tree_fig9_image_cs.png", 0.3, 1.4, w=7.5)
bl(s, 8.2, 1.6, 4.5, 5, [
    "Real wavelet coefficients from",
    "8 natural images (Haar DWT)",
    "",
    "TH-hybrid achieves:",
    "  \u2022 Best SSIM at all CS ratios",
    "  \u2022 31.66 dB PSNR at CS 50%",
    "  \u2022 Beats HyperLISTA by 10+ dB",
    "",
    "Tree structure in wavelet domain",
    "\u2192 tree-aware proximal outperforms",
    "   elementwise proximal",
], 18)

# === SLIDE 11: Section header ===
s = prs.slides.add_slide(prs.slide_layouts[6])
sec(s, "Appendix: Extended Analysis")

# === SLIDE 12: Backbone Ablation (App B) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix B: Backbone Ablation")
img(s, "tree_fig_backbone.png", 0.3, 1.4, w=6.5)
bl(s, 7.3, 1.6, 5.5, 5.5, [
    "6 backbone variants compared:",
    "",
    "Tree-ALISTA (32p):       \u221215.00 dB",
    "Tree-ALISTA-MM (32p):    \u221217.02 dB",
    "ALISTA-MM-Sym (32p):     +38.80 (diverges!)",
    "TH-Elem (3p, no tree):   \u221215.75 dB",
    "Tree-HyperLISTA (3p):    \u221228.66 dB",
    "",
    "Key: tree-aware operator is the",
    "critical ingredient, not backbone",
    "complexity or parameter count",
], 18)

# === SLIDE 13: Extended Mechanisms (App C) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix C: Extended Mechanism Ablation")
img(s, "tree_fig_mechanisms_ext.png", 0.3, 1.4, w=6)
bl(s, 7, 1.6, 5.5, 5, [
    "Differentiable relaxations tested:",
    "",
    "Tree Hard (M1):        \u22123.62 dB",
    "Diff. Ancestor:        \u22123.12 dB",
    "Gumbel Top-K:          \u22122.73 dB",
    "Threshold (M2):        \u22122.73 dB",
    "Hybrid (M3):           \u22122.73 dB",
    "",
    "Non-differentiable hard projection",
    "outperforms smooth relaxations",
], 18)

# === SLIDE 14: Rho + Sparsity (App D, E) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix D-E: Rho Sensitivity & Sparsity Mismatch")
img(s, "tree_fig_rho_dense.png", 0.3, 1.4, w=5.5)
img(s, "tree_fig_sparsity.png", 6.5, 1.4, w=6)
bl(s, 0.5, 5.8, 12, 1.2, [
    "Left: \u03c1=0.1 and \u03c1=0.5 give same result \u2014 robust to decay choice  |  Right: graceful at sparser; all methods degrade when denser",
], 16, c=GRAY)

# === SLIDE 15: Extended Mismatch (App F) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix F: Extended Mismatch Robustness")
img(s, "tree_fig_mm_magnitude.png", 0.2, 1.3, w=6.3)
img(s, "tree_fig_mm_consistency.png", 6.8, 1.3, w=6)
tb(s, 0.5, 5.6, 12, 1, "Magnitude scaling: TH robust (\u221228.7\u2192\u221228.1 dB at 3\u00d7) while HyperLISTA diverges (+85 dB)", 16, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.5, 6.1, 12, 1, "Tree consistency: TH degrades from \u221228.8 to \u22127.2 dB as tree structure breaks \u2014 expected behavior", 16, c=GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 16: Low-SNR + Superlinear (App G, H) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix G-H: Low-SNR & Superlinear Convergence")
img(s, "tree_fig_lowsnr.png", 0.3, 1.4, w=5.8)
img(s, "tree_fig_superlinear.png", 6.8, 1.4, w=5.8)
bl(s, 0.5, 5.8, 12, 1.2, [
    "Left: TH maintains advantage at all SNR levels (10/20/30 dB)  |  Right: per-instance tuning \u2192 \u22124.0 dB at only 6 layers (future work)",
], 16, c=GRAY)

# === SLIDE 17: Cross-Structure (App I) ===
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, "Appendix I: Cross-Structure Experiments")
img(s, "tree_fig_cross.png", 1, 1.4, w=6)
bl(s, 7.5, 1.6, 5, 5, [
    "Train on tree-sparse, test on elementwise:",
    "",
    "TH-hybrid on non-tree data: \u22123.45 dB",
    "TH-hybrid on tree data:     \u22122.73 dB",
    "",
    "Tree-aware method does NOT hurt",
    "performance on non-tree signals",
    "",
    "\u2192 Safe to use tree operator even",
    "   when structure is uncertain",
], 18)

# === SLIDE 18: Key Takeaways ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_BG)
tb(s, 0.8, 0.4, 11, 0.8, "Key Takeaways", 40, b=True, c=BLUE, a=PP_ALIGN.CENTER)
takeaways = [
    ("1", "Structure > Parameters", "3 params + tree operator beats 1.56M params (LISTA) by 15.5 dB"),
    ("2", "Minimal Parameterization Generalizes", "3-parameter principle extends from elementwise to tree-structured sparsity"),
    ("3", "Hard Projection Works", "Simple greedy top-K tree projection nearly matches hybrid mode"),
    ("4", "Real-World Validation", "Best SSIM on wavelet image CS at all compression ratios"),
    ("5", "Robust & Safe", "Graceful degradation under mismatch; doesn't hurt non-tree signals"),
]
for i, (num, title, desc) in enumerate(takeaways):
    y = 1.3 + i * 1.15
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y+0.05), Inches(0.6), Inches(0.6))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()
    tf = sh.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb(s, 2.2, y, 9, 0.5, title, 24, b=True, c=DARK)
    tb(s, 2.2, y+0.45, 9, 0.5, desc, 17, c=GRAY)

# === SLIDE 19: Thank You ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BLUE)
tb(s, 1, 1.5, 11.3, 1.5, "Thank You", 52, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 1, 3.2, 11.3, 0.8, "Tree-HyperLISTA: 3 parameters, 12.8 dB better than HyperLISTA, perfect recall", 24, c=RGBColor(0xBB,0xDD,0xFF), a=PP_ALIGN.CENTER)
tb(s, 1, 4.5, 11.3, 0.8, "Code: github.com/ALPHAGOD12/tree-hyperlista", 20, c=RGBColor(0x99,0xCC,0xFF), a=PP_ALIGN.CENTER)
tb(s, 1, 5.3, 11.3, 0.8, "Shresth Verma & Aditya Agrawal  \u2022  IIT Bombay", 20, c=RGBColor(0x99,0xCC,0xFF), a=PP_ALIGN.CENTER)

out = "tree_hyperlista_slides.pptx"
prs.save(out)
print(f"Saved {out} ({len(prs.slides)} slides)")
